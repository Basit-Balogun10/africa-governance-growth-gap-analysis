"""
Fix PowerPoint presentation:
1. Remove \n\n line breaks on specific slides
2. Add submission links to first and last slides
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# Load presentation
prs = Presentation(r'C:\Users\USER\Documents\10Alatytics-2025\presentation\10Alytics_Governance_Growth_Gap.pptx')

print(f"Total slides: {len(prs.slides)}\n")

# Fix slide text replacements
fixes = {
    # Find and replace \n\n with proper line breaks
    'We analyzed 14 African countries across 65 years\\n\\nand found the OPPOSITE.': 
        'We analyzed 14 African countries across 65 years\nand found the OPPOSITE.',
    'Questions? Challenges?\\n\\nLet\'s Discuss.': 
        'Questions? Challenges?\nLet\'s Discuss.',
}

# Scan and fix slides
for slide_num, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for bad_text, good_text in fixes.items():
                        if bad_text in run.text:
                            print(f"✓ Fixed on slide {slide_num}: {bad_text[:50]}...")
                            run.text = run.text.replace(bad_text, good_text)

# Add submission links to slide 1 (after title)
slide_1 = prs.slides[0]
# Add text box at bottom
left = Inches(0.5)
top = Inches(6.5)
width = Inches(9)
height = Inches(0.8)

textbox = slide_1.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "📦 Complete Submission: GitHub (Live Code) + ZIP (All Files - Extract to View)"
p.font.size = Pt(10)
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

p2 = text_frame.add_paragraph()
p2.text = "https://github.com/Basit-Balogun10/africa-governance-growth-gap-analysis"
p2.font.size = Pt(9)
p2.font.bold = True
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER

print(f"\n✓ Added submission links to slide 1")

# Add to last slide as well
last_slide = prs.slides[-1]
textbox_last = last_slide.shapes.add_textbox(left, Inches(6), width, height)
text_frame_last = textbox_last.text_frame
text_frame_last.word_wrap = True

p_last = text_frame_last.paragraphs[0]
p_last.text = "📦 Full Submission ZIP (13MB - Please Extract): See repo or contact basitbalogun10@gmail.com"
p_last.font.size = Pt(10)
p_last.font.name = 'Calibri'
p_last.alignment = PP_ALIGN.CENTER

p_last2 = text_frame_last.add_paragraph()
p_last2.text = "GitHub: https://github.com/Basit-Balogun10/africa-governance-growth-gap-analysis"
p_last2.font.size = Pt(9)
p_last2.font.bold = True
p_last2.font.name = 'Calibri'
p_last2.alignment = PP_ALIGN.CENTER

print(f"✓ Added submission links to slide {len(prs.slides)}")

# Save
prs.save(r'C:\Users\USER\Documents\10Alatytics-2025\presentation\10Alytics_Governance_Growth_Gap.pptx')
print(f"\n✅ PowerPoint updated successfully!")
